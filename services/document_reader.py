"""
=========================================================
Document Reader Service
=========================================================
Membaca berbagai format dokumen menjadi teks.
Menghasilkan object Dokumen dan MataKuliah.
=========================================================
"""

from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from services.models import Dokumen, MataKuliah


class DocumentReader:

    SUPPORTED_FORMATS = (
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
    )

    def __init__(self):

        self.file_dibaca = 0
        self.file_gagal = 0
        self.total_karakter = 0

    # =====================================================
    # FORMAT
    # =====================================================

    def format_didukung(self, file_path: Path):

        return file_path.suffix.lower() in self.SUPPORTED_FORMATS

    # =====================================================
    # PDF
    # =====================================================

    def baca_pdf(self, file_path: Path):

        reader = PdfReader(file_path)

        teks = []

        for halaman in reader.pages:

            isi = halaman.extract_text()

            if isi:
                teks.append(isi)

        return "\n".join(teks)

    # =====================================================
    # DOCX
    # =====================================================

    def baca_docx(self, file_path: Path):

        document = Document(file_path)

        teks = []

        for paragraph in document.paragraphs:

            if paragraph.text.strip():
                teks.append(paragraph.text)

        for table in document.tables:

            for row in table.rows:

                isi = []

                for cell in row.cells:

                    if cell.text.strip():
                        isi.append(cell.text.strip())

                if isi:
                    teks.append(" | ".join(isi))

        return "\n".join(teks)

    # =====================================================
    # PPTX
    # =====================================================

    def baca_pptx(self, file_path: Path):

        presentation = Presentation(file_path)

        teks = []

        for nomor, slide in enumerate(
            presentation.slides,
            start=1
        ):

            teks.append(
                f"--- SLIDE {nomor} ---"
            )

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():
                        teks.append(shape.text)

        return "\n".join(teks)

    # =====================================================
    # XLSX
    # =====================================================

    def baca_xlsx(self, file_path: Path):

        workbook = load_workbook(
            file_path,
            data_only=True
        )

        teks = []

        for worksheet in workbook.worksheets:

            teks.append(
                f"--- SHEET: {worksheet.title} ---"
            )

            for row in worksheet.iter_rows(
                values_only=True
            ):

                nilai = []

                for cell in row:

                    if cell is not None:
                        nilai.append(str(cell))

                if nilai:
                    teks.append(" | ".join(nilai))

        return "\n".join(teks)

    # =====================================================
    # SATU FILE
    # =====================================================

    def baca_file(self, file_path: Path):

        try:

            suffix = file_path.suffix.lower()

            if suffix == ".pdf":

                isi = self.baca_pdf(file_path)

            elif suffix == ".docx":

                isi = self.baca_docx(file_path)

            elif suffix == ".pptx":

                isi = self.baca_pptx(file_path)

            elif suffix == ".xlsx":

                isi = self.baca_xlsx(file_path)

            else:

                return None

            self.file_dibaca += 1
            self.total_karakter += len(isi)

            return Dokumen(

                nama=file_path.name,

                path=str(file_path),

                kategori=file_path.parent.name,

                format=suffix,

                ukuran=file_path.stat().st_size,

                isi=isi

            )

        except Exception:

            self.file_gagal += 1

            return None

    # =====================================================
    # SATU FOLDER
    # =====================================================

    def baca_folder(self, folder_path: Path):

        hasil = []

        if not folder_path.exists():

            return hasil

        for file in sorted(folder_path.iterdir()):

            if not file.is_file():

                continue

            dokumen = self.baca_file(file)

            if dokumen:

                hasil.append(dokumen)

        return hasil

    # =====================================================
    # SATU MATA KULIAH
    # =====================================================

    def baca_mata_kuliah(self, folder_mk: Path):

        modul = self.baca_folder(
            folder_mk / "modul"
        )

        artefak = self.baca_folder(
            folder_mk / "artefak"
        )

        jumlah_file = len(modul) + len(artefak)

        jumlah_karakter = sum(
            len(doc.isi)
            for doc in modul + artefak
        )

        return MataKuliah(

            nama=folder_mk.name,

            modul=modul,

            artefak=artefak,

            jumlah_file=jumlah_file,

            jumlah_karakter=jumlah_karakter

        )

    # =====================================================
    # RINGKASAN
    # =====================================================

    def ringkasan(self):

        return {

            "file_dibaca": self.file_dibaca,

            "file_gagal": self.file_gagal,

            "total_karakter": self.total_karakter

        }